"""Build: AI Integration Plan deck"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy
import copy
from pathlib import Path

PRESENTATION_DATE = '260508'
PRESENTATION_NAME = 'AI_Integration_Plan'

def find_workspace_root(start):
    for path in [start, *start.parents]:
        if (path / '_templates' / 'open-template.pptx').exists():
            return path
    raise FileNotFoundError('Could not find _templates/open-template.pptx')

SCRIPT_DIR = Path(__file__).resolve().parent
WORKSPACE_ROOT = find_workspace_root(SCRIPT_DIR)
TEMPLATE = WORKSPACE_ROOT / '_templates' / 'open-template.pptx'
OUTPUT = SCRIPT_DIR / 'output' / f'{PRESENTATION_DATE}_{PRESENTATION_NAME}.pptx'

# Sample template colours (open template)
DARK       = RGBColor(0x0D, 0x0D, 0x0D)
GREY       = RGBColor(0x99, 0x99, 0x99)
ACCENT     = RGBColor(0x00, 0x30, 0x57)   # dark navy — template primary
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF2, 0xF2, 0xF2)

prs = Presentation(str(TEMPLATE))
# Remove template example slides — keep masters only
slide_ids = list(prs.slides._sldIdLst)
for sld_id in slide_ids:
    rId = sld_id.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sld_id)

LAYOUT_COVER   = prs.slide_layouts[0]   # Cover Title
LAYOUT_CONTENT = prs.slide_layouts[2]   # 1_Title & Contents
LAYOUT_CHAPTER = prs.slide_layouts[4]   # Chapter Title
LAYOUT_BACK    = prs.slide_layouts[5]   # Back Cover

# ── helpers ──────────────────────────────────────────────────────────────────

def add_slide(layout):
    return prs.slides.add_slide(layout)

def set_text(shape, text, size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or DARK

def add_textbox(slide, left, top, width, height, text, size=12, bold=False,
                color=None, align=PP_ALIGN.LEFT, mono=False, bg=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    if bg:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg

    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or DARK
        if mono:
            run.font.name = 'Courier New'
    return txBox

def title_and_body(layout, title_text, body_text, body_size=13):
    slide = add_slide(layout)
    slide.placeholders[0].text = title_text
    body = slide.placeholders[10] if 10 in [ph.placeholder_format.idx
                                             for ph in slide.placeholders] else None
    if body and body_text:
        tf = body.text_frame
        tf.word_wrap = True
        tf.text = ''
        for i, line in enumerate(body_text.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(body_size)
            run.font.color.rgb = DARK
    return slide

# ── SLIDE 1: Title Cover ──────────────────────────────────────────────────────
slide = add_slide(LAYOUT_COVER)
slide.placeholders[0].text = 'AI Integration Plan'
slide.placeholders[1].text = 'Presenter Name  |  8 May 2026'

# ── SLIDE 2: AI Is Already Here ───────────────────────────────────────────────
slide = title_and_body(LAYOUT_CONTENT, 'AI Is Already Here', '')
add_textbox(slide, 0.67, 1.55, 11.6, 4.5,
    'AI is already being used in consulting today — not as a future concept, but as a working layer in how firms operate.\n\n'
    'Professional services firms are using AI to handle writing, research, structured outputs, and client communications.\n\n'
    'The question is no longer "will this happen?"\n\n'
    'It already has.\n\n'
    'The question is: how should this organization adopt AI — and what do we want it to do?',
    size=14)

# ── SLIDE 3: The Wrong Question ───────────────────────────────────────────────
slide = add_slide(LAYOUT_CONTENT)
slide.placeholders[0].text = 'The Wrong Question'

add_textbox(slide, 0.67, 1.55, 11.6, 1.2,
    '"Will AI replace consultants?"',
    size=18, bold=False, color=GREY, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.67, 2.8, 11.6, 0.5,
    '── That is the wrong question. ──',
    size=13, color=GREY, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.67, 3.4, 11.6, 0.8,
    'The right question:',
    size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.67, 4.0, 11.6, 1.2,
    '"Which parts of consulting work should AI handle,\nand which parts should stay with the consultant?"',
    size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.67, 5.4, 11.6, 0.7,
    'The highest-value parts of consulting — reading a client, diagnosing the real problem,\n'
    'building trust, making judgment calls — cannot be automated.',
    size=12, color=DARK, align=PP_ALIGN.CENTER)

# ── SLIDE 4: Two Layers of Consulting Work ────────────────────────────────────
slide = add_slide(LAYOUT_CONTENT)
slide.placeholders[0].text = 'Two Layers of Consulting Work'

# Left column — operating layer
add_textbox(slide, 0.67, 1.55, 5.6, 0.55,
    'OPERATING LAYER', size=13, bold=True, color=GREY)
add_textbox(slide, 0.67, 2.1, 5.6, 0.35,
    'AI assists in this layer', size=11, color=GREY)
add_textbox(slide, 0.67, 2.5, 5.6, 3.0,
    'Setting up workspaces\nLogging meeting notes\nFormatting deliverables\nRouting outputs to the right place\nDrafting status updates',
    size=13, mono=True)

# Right column — judgment layer
add_textbox(slide, 6.8, 1.55, 5.8, 0.55,
    'JUDGMENT LAYER', size=13, bold=True, color=ACCENT)
add_textbox(slide, 6.8, 2.1, 5.8, 0.35,
    'The consultant keeps this', size=11, color=ACCENT)
add_textbox(slide, 6.8, 2.5, 5.8, 3.0,
    'Diagnosing the real problem\nReading client dynamics\nMaking recommendations\nBuilding trust\nDeciding what matters',
    size=13, mono=True, color=DARK)

add_textbox(slide, 0.67, 5.7, 11.6, 0.5,
    'The AI integration plan is built on this distinction.',
    size=12, color=GREY, align=PP_ALIGN.CENTER)

# ── SLIDE 5: Where AI Fits in Your Workflow ───────────────────────────────────
slide = add_slide(LAYOUT_CONTENT)
slide.placeholders[0].text = 'Where AI Fits in Your Workflow'

diagram = """\
Intake +         Scope +          Delivery         Review +          Handoff
Discovery    ->  Proposal     ->  Workspace    ->  Client Comms  ->  + Case Study
     |                |                |                 |                 |
  [Brief           [Scope &        [Formatted        [Status           [Handoff
   Builder]         Risk Flags]     Outputs]          Updates]          Package]
     |                |                |                 |                 |
  Structures       Identifies      Routes to         Drafts client     Packages
  messy input      scope gaps      right folder      updates & logs    docs + files
                                                                      ___________
<----------------------- AI assists in this layer --------------------------------->"""

add_textbox(slide, 0.47, 1.55, 12.4, 3.2, diagram,
    size=10, mono=True, bg=LIGHT_BG)

add_textbox(slide, 0.47, 4.85, 12.4, 0.45,
    'You still run every stage. AI assists with the parts that are repetitive, '
    'forgettable, or time-consuming.',
    size=11, color=GREY)

# ── SLIDE 6: What the Tools Look Like ────────────────────────────────────────
slide = add_slide(LAYOUT_CONTENT)
slide.placeholders[0].text = 'What the Tools Look Like'

toolkit = """\
+----------------------------------------------------------+
|          CONSULTANT AI TOOLKIT — PHASE 1                 |
+----------------------------------------------------------+
| 1. Intake Brief Builder                                  |
|    Turns a messy client email or transcript into a       |
|    structured brief: problem, assumptions, risks,        |
|    and missing information.                              |
| 2. Meeting Notes Processor                               |
|    Turns a transcript into decisions, action items,      |
|    risks, and follow-ups. Saves in the right place.      |
| 3. Review Gate                                           |
|    Checks a deliverable before it leaves your desk:      |
|    scope match, unsupported claims, formatting.          |
| 4. Handoff Package Builder                               |
|    Creates documentation so clients can use what you     |
|    built — not just "send the files and hope."           |
+----------------------------------------------------------+"""

add_textbox(slide, 0.47, 1.55, 12.4, 3.8, toolkit,
    size=10, mono=True, bg=LIGHT_BG)

add_textbox(slide, 0.47, 5.45, 12.4, 0.45,
    'Each tool does one job. You use it at the right moment in the workflow you already follow.',
    size=11, color=GREY)

# ── SLIDE 7: We've Already Started ────────────────────────────────────────────
slide = add_slide(LAYOUT_CONTENT)
slide.placeholders[0].text = "We've Already Started"

add_textbox(slide, 0.67, 1.7, 11.6, 0.8,
    'A sub-trial of a single-use consultant toolkit is currently underway inside the organization.',
    size=16, bold=True)

add_textbox(slide, 0.67, 2.7, 11.6, 3.0,
    'Focus area:    sample implementation projects\n\n'
    'Participants:  Trial Participant A  |  Consultant\n'
    '               Trial Participant B  |  Consultant\n\n'
    'Status:        Active — running on live engagements',
    size=13, mono=True, bg=LIGHT_BG)

add_textbox(slide, 0.67, 6.0, 11.6, 0.6,
    'We will hear directly from the trial participants at the end of this session.',
    size=13, bold=True, color=ACCENT)

# ── SLIDE 8: The Plan ─────────────────────────────────────────────────────────
slide = add_slide(LAYOUT_CONTENT)
slide.placeholders[0].text = 'The Plan'

plan = """\
Phase 1: Core Delivery           Phase 2: Control + Memory    Phase 3: Growth [LATER]
[IN TRIAL — sample implementation toolkit]       [NEXT]                       ----------------------
Workspace Generator              Change Request Detector       Proposal Builder
Intake Brief Builder             Workspace Auditor             Case Study Anonymizer
Scope Agreement Builder          Decision Log Builder          Deliverable Builder+
Meeting Notes Processor          Consultant Command Center     Practice Analytics
Review Gate
Client Comms Drafter
Handoff Package Builder"""

add_textbox(slide, 0.47, 1.55, 12.4, 2.8, plan,
    size=10, mono=True, bg=LIGHT_BG)

add_textbox(slide, 0.47, 4.45, 12.4, 0.45,
    'Phase 1 is live. Phases 2 and 3 follow once the core delivery loop is proven.',
    size=12, color=GREY)

# ── SLIDE 9: What This Means for You ──────────────────────────────────────────
slide = add_slide(LAYOUT_CONTENT)
slide.placeholders[0].text = 'What This Means for You'

add_textbox(slide, 0.67, 1.7, 11.6, 0.7,
    'You do not need to change how you think about consulting.',
    size=18, bold=True)

add_textbox(slide, 0.67, 2.6, 11.6, 3.5,
    'You need to start using the tools where they fit.\n\n'
    '→  Each tool does one job. You do not need to learn a new system.\n\n'
    '→  Start with whichever tool addresses the most friction in your current work.\n\n'
    '→  The consultants who adopt early will have more capacity for the\n'
    '   work that only they can do.\n\n'
    '→  More detail on rollout and onboarding to follow.',
    size=13)

# ── SLIDE 10: Live Demo ───────────────────────────────────────────────────────
slide = add_slide(LAYOUT_CONTENT)
slide.placeholders[0].text = 'Live Demo: Deck-Making Toolkit in Action'

add_textbox(slide, 0.67, 1.55, 11.6, 0.55,
    'This deck was built using the AI integration toolkit — end to end.',
    size=14, bold=True)

demo = """\
+------------------------------------------------------------+
|  DECK BUILDER TOOLKIT — WHAT JUST HAPPENED                |
+------------------------------------------------------------+
|  Brief       ->  Source Discovery  ->  Narrative           |
|  Storyboard  ->  Design Plan       ->  Build (this deck)   |
+------------------------------------------------------------+
|  Each stage: discuss with user -> agree -> produce output  |
|  All outputs saved, versioned, and reusable                |
+------------------------------------------------------------+"""

add_textbox(slide, 0.47, 2.25, 12.4, 2.5, demo,
    size=10, mono=True, bg=LIGHT_BG)

add_textbox(slide, 0.67, 4.85, 11.6, 0.5,
    'The same pattern applies to every consultant workflow — not just decks.',
    size=13, color=GREY)

add_textbox(slide, 0.67, 5.45, 11.6, 0.9,
    'What you just watched:\n'
    'Brief → sources → argument → slides → design → .pptx — in one session, with full audit trail.',
    size=12, color=DARK)

# ── SLIDE 11: Over to the Trial Team ──────────────────────────────────────────
slide = add_slide(LAYOUT_CONTENT)
slide.placeholders[0].text = 'The Trial Team'

add_textbox(slide, 0.67, 1.55, 11.6, 0.6,
    'Two consultants have been running the toolkit on live sample implementation engagements.',
    size=13, color=GREY)

# Trial Participant A
add_textbox(slide, 1.0, 2.4, 5.0, 0.6,
    '[ Photo ]', size=12, color=GREY, align=PP_ALIGN.CENTER, bg=LIGHT_BG)
add_textbox(slide, 1.0, 3.1, 5.0, 0.5,
    'Trial Participant A', size=16, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, 1.0, 3.7, 5.0, 0.4,
    'Consultant', size=12, color=GREY, align=PP_ALIGN.CENTER)

# Trial Participant B
add_textbox(slide, 7.3, 2.4, 5.0, 0.6,
    '[ Photo ]', size=12, color=GREY, align=PP_ALIGN.CENTER, bg=LIGHT_BG)
add_textbox(slide, 7.3, 3.1, 5.0, 0.5,
    'Trial Participant B', size=16, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, 7.3, 3.7, 5.0, 0.4,
    'Consultant', size=12, color=GREY, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.67, 5.8, 11.6, 0.6,
    '"The goal is not to automate the consultant out of the process. The goal is to remove '
    'repeated setup, forgotten context, inconsistent review, and weak handoff — so the '
    'consultant can spend more time on judgment, diagnosis, and client trust."',
    size=11, color=GREY, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
import os
os.makedirs(OUTPUT.parent, exist_ok=True)
prs.save(str(OUTPUT))
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
