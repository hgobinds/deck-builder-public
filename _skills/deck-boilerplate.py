"""
Deck Builder — Boilerplate
Copy into 06_build/build_deck.py for each new deck project.
Rename OUTPUT to match YYMMDD_Presentation_Name.pptx using presentation date.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

PRESENTATION_DATE = 'YYMMDD'
PRESENTATION_NAME = 'Presentation_Name'

def find_workspace_root(start):
    for path in [start, *start.parents]:
        if (path / '_templates' / 'open-template.pptx').exists():
            return path
    raise FileNotFoundError('Could not find _templates/open-template.pptx')

SCRIPT_DIR = Path(__file__).resolve().parent
WORKSPACE_ROOT = find_workspace_root(SCRIPT_DIR)
TEMPLATE = WORKSPACE_ROOT / '_templates' / 'open-template.pptx'
OUTPUT = SCRIPT_DIR / 'output' / f'{PRESENTATION_DATE}_{PRESENTATION_NAME}.pptx'

# Sample template colours
DARK     = RGBColor(0x0D, 0x0D, 0x0D)
GREY     = RGBColor(0x99, 0x99, 0x99)
ACCENT   = RGBColor(0x00, 0x30, 0x57)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF2, 0xF2, 0xF2)

prs = Presentation(str(TEMPLATE))

# Remove template example slides
slide_ids = list(prs.slides._sldIdLst)
for sld_id in slide_ids:
    rId = sld_id.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sld_id)

# Layout index
# [0] Cover Title    [1] 2_Title & Contents    [2] 1_Title & Contents (default)
# [3] 3_Title & Contents    [4] Chapter Title    [5] Back Cover
LAYOUT_COVER   = prs.slide_layouts[0]
LAYOUT_CONTENT = prs.slide_layouts[2]
LAYOUT_CHAPTER = prs.slide_layouts[4]
LAYOUT_BACK    = prs.slide_layouts[5]

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide(layout):
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text,
                size=12, bold=False, color=None, align=PP_ALIGN.LEFT,
                mono=False, bg=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if bg:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or DARK
        if mono:
            run.font.name = 'Courier New'
    return txBox

# ── Slides ────────────────────────────────────────────────────────────────────
# Add slides here following the approved storyboard.
# Standard content slide:
#   slide = add_slide(LAYOUT_CONTENT)
#   slide.placeholders[0].text = 'Slide Title'
#   add_textbox(slide, 0.67, 1.55, 11.6, 4.5, 'Body text...', size=13)
#
# ASCII diagram slide:
#   add_textbox(slide, 0.47, 1.55, 12.4, <height>, diagram, size=10, mono=True, bg=LIGHT_BG)
#   Set height to fit content — do not leave large empty grey areas.
#
# Two-column slide (left / right):
#   add_textbox(slide, 0.67, 1.55, 5.6, 3.5, left_text, ...)
#   add_textbox(slide, 6.8,  1.55, 5.8, 3.5, right_text, ...)

# ── Save ──────────────────────────────────────────────────────────────────────
import os
os.makedirs(OUTPUT.parent, exist_ok=True)
prs.save(str(OUTPUT))
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
